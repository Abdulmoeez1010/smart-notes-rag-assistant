from pptx import Presentation
from langchain_core.documents import Document
from app.ingestion.pdf_loader import generate_doc_id


def load_pptx(file_path: str, original_filename: str = "presentation.pptx") -> Document:
    prs = Presentation(file_path)
    title = get_pptx_title(prs, original_filename)
    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
        if texts:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
    if not slides_text:
        raise ValueError(f"No text extracted from PPTX: {file_path}")
    full_text = "\n\n".join(slides_text)
    doc_id = generate_doc_id(file_path)
    return Document(page_content=full_text, metadata={"doc_id": doc_id, "source_type": "pptx", "title": title})


def get_pptx_title(prs, original_filename: str) -> str:
    """Try the first slide's title placeholder, fall back to filename."""
    try:
        first_slide = prs.slides[0]
        if first_slide.shapes.title and first_slide.shapes.title.text.strip():
            return first_slide.shapes.title.text.strip()
    except Exception:
        pass
    return original_filename